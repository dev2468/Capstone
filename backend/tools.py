"""
tools.py — Shared MCP tool implementations.

Both api.py (FastAPI/async) and server.py (stdio_server/async) import from
this module so that security logic, file-parsing, and command execution stay
in one place.

Public surface
--------------
read_any_file(filepath)                         → str
list_files_impl(folder)                         → str
run_command_impl(command, *, timeout)           → str   (async)
read_file_impl(filepath)                        → str

Security helpers (re-exported so callers can validate before calling)
----------------------------------------------------------------------
is_command_allowed(command)                     → tuple[bool, str]
DANGEROUS_PATTERNS                              re.Pattern
ALLOWED_COMMAND_PREFIXES                        set[str]
"""

import asyncio
import os
import re
import logging

log = logging.getLogger(__name__)

# ── Security constants ────────────────────────────────────────────────────────
DANGEROUS_PATTERNS = re.compile(r'[;&`<>|]|\$[\(\{]')
ALLOWED_COMMAND_PREFIXES: set[str] = {
    "git", "python", "pip", "uv", "npm", "node", "gradlew",
    "java", "adb", "dir", "type", "echo", "curl",
    "get-childitem", "get-content", "test-path"
}


def is_command_allowed(command: str) -> tuple[bool, str]:
    """Return (True, '') if *command* passes security checks, else (False, reason)."""
    if DANGEROUS_PATTERNS.search(command):
        return False, "Command contains forbidden character/sequence"
    parts = command.strip().split()
    if not parts:
        return False, "Empty command"
    prefix = parts[0].strip("'\"")
    if prefix.lower() not in ALLOWED_COMMAND_PREFIXES:
        return False, f"Command prefix '{prefix}' is not in the allowlist."
    return True, ""


# ── File reader ───────────────────────────────────────────────────────────────
def read_any_file(filepath: str) -> str:
    """Read *filepath* and return its contents as a string.

    Handles plain text, PDF, DOCX, XLSX, PPTX, and common image types.
    Falls back to UTF-8 text for any unrecognised extension.
    """
    ext = os.path.splitext(filepath)[1].lower()

    if ext in {'.txt', '.py', '.js', '.ts', '.json', '.csv', '.md',
               '.html', '.css', '.xml'}:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    elif ext == '.pdf':
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        return '\n'.join([page.extract_text() or '' for page in reader.pages])

    elif ext == '.docx':
        from docx import Document
        doc = Document(filepath)
        return '\n'.join([para.text for para in doc.paragraphs])

    elif ext == '.xlsx':
        from openpyxl import load_workbook
        wb = load_workbook(filepath, data_only=True)
        lines = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            lines.append(f"Sheet: {sheet}")
            for row in ws.iter_rows(values_only=True):
                lines.append(str(row))
        return '\n'.join(lines)

    elif ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(filepath)
        lines = []
        for i, slide in enumerate(prs.slides):
            lines.append(f"Slide {i + 1}:")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    lines.append(shape.text)
        return '\n'.join(lines)

    elif ext in {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'}:
        from PIL import Image
        img = Image.open(filepath)
        return f"Image: {img.format}, Size: {img.size[0]}x{img.size[1]}px, Mode: {img.mode}"

    else:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()


# ── Shared tool implementations ───────────────────────────────────────────────

def list_files_impl(folder: str) -> str:
    """Return a newline-separated listing of *folder*, or an empty-folder message.

    Raises:
        ValueError: if *folder* is empty.
        FileNotFoundError: if *folder* does not exist.
    """
    folder = folder.strip()
    if not folder:
        raise ValueError("No folder path provided")
    if not os.path.exists(folder):
        raise FileNotFoundError(f"Folder does not exist: {folder}")
    files = os.listdir(folder)
    if not files:
        return "Folder is empty"
    if len(files) > 40:
        more = len(files) - 40
        return '\n'.join(files[:40]) + f"\n... {more} more items truncated, narrow your folder path for full listing."
    return '\n'.join(files)


async def run_command_impl(command: str, *, timeout: float) -> str:
    """Run *command* in PowerShell and return combined stdout+stderr.

    Uses ``asyncio.create_subprocess_exec`` (non-blocking) and
    ``asyncio.wait_for`` to enforce *timeout* without blocking the event loop.

    Raises:
        ValueError: if *command* is empty or blocked by the allowlist.
        asyncio.TimeoutError: if the process does not finish within *timeout* seconds.
        PermissionError: propagated from the OS.
    """
    command = command.strip()
    if not command:
        raise ValueError("No command provided")

    # Defensive backstop: strip unnecessary powershell wrapper if the model provides one
    wrapper_match = re.match(r'^powershell(?:\.exe)?\s+-Command\s+(.*)', command, re.IGNORECASE)
    if wrapper_match:
        log.warning(f"Stripping unnecessary powershell wrapper from command: {command}")
        command = wrapper_match.group(1).strip()
        # Remove surrounding quotes if they enclose the whole remaining command
        if (command.startswith('"') and command.endswith('"')) or (command.startswith("'") and command.endswith("'")):
            command = command[1:-1].strip()

    if not command:
        raise ValueError("No command provided")

    allowed, reason = is_command_allowed(command)
    if not allowed:
        raise ValueError(f"Command blocked by security policy. Reason: {reason}")

    env = os.environ.copy()
    env["PATH"] = r"C:\Users\HP\Desktop\MCP\venv\Scripts" + ";" + env.get("PATH", "")
    proc = await asyncio.create_subprocess_exec(
        "powershell", "-ExecutionPolicy", "RemoteSigned",
        "-NonInteractive", "-NoProfile", "-Command", command,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
        env=env,
    )
    try:
        stdout_b, stderr_b = await asyncio.wait_for(
            proc.communicate(), timeout=timeout
        )
    except asyncio.TimeoutError:
        proc.kill()
        await proc.communicate()   # reap zombie process
        raise

    stdout_str = stdout_b.decode(errors="replace")
    stderr_str = stderr_b.decode(errors="replace")

    if len(stdout_str) > 20000:
        stdout_str = stdout_str[:20000] + "\n... output truncated at 20000 chars."
    if len(stderr_str) > 20000:
        stderr_str = stderr_str[:20000] + "\n... output truncated at 20000 chars."

    output = stdout_str + stderr_str
    return output.strip() if output.strip() else "Command ran with no output"


def read_file_impl(filepath: str, start_line: int | None = None, end_line: int | None = None) -> str:
    """Read and return the contents of *filepath*.

    Args:
        filepath:   Path to the file.
        start_line: 1-based line number to start from (inclusive). Optional.
        end_line:   1-based line number to stop at (inclusive). Optional.

    Returns the full file up to 50000 chars. For larger files, use
    start_line/end_line to read in chunks.
    """
    filepath = filepath.strip()
    if not filepath:
        raise ValueError("No filepath provided")
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File does not exist: {filepath}")

    content = read_any_file(filepath)

    if start_line is not None or end_line is not None:
        lines = content.splitlines(keepends=True)
        total = len(lines)
        s = (start_line - 1) if start_line else 0
        e = end_line if end_line else total
        s = max(0, s)
        e = min(total, e)
        content = f"[Lines {s+1}–{e} of {total}]\n" + "".join(lines[s:e])

    content = content.strip()
    if not content:
        return "File is empty"

    if len(content) > 50000:
        return content[:50000] + "\n... truncated at 50000 chars. Use start_line/end_line to read further."

    return content


ALLOWED_WRITE_PATHS = [
    "C:/DevMCP",
    "C:/Users/HP/Desktop/MCP",
    "C:/Users/HP/Downloads"
]

def write_file_impl(filepath: str, content: str) -> str:
    """Write content to filepath, with path restriction enforcement."""
    # Normalize to forward slashes for consistent comparison
    normalized = os.path.normpath(filepath).replace("\\", "/")
    if not any(normalized.startswith(os.path.normpath(p).replace("\\", "/"))
               for p in ALLOWED_WRITE_PATHS):
        raise ValueError(
            f"Write blocked: '{filepath}' is outside allowed paths: "
            f"{ALLOWED_WRITE_PATHS}"
        )
    os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)
    tmp = filepath + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        f.write(content)
        f.flush()
        os.fsync(f.fileno())
    os.replace(tmp, filepath)
    return f"Successfully wrote {len(content)} characters to {filepath}"
